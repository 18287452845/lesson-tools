import pathlib

import pytest
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.enum import text as pptx_text_enum

from backend.main import app
from backend.services.builtin_template import (
    BUILTIN_TEMPLATE_ID,
    REQUIRED_TEMPLATE_FIELDS,
    get_builtin_template_path,
    require_valid_builtin_template,
    validate_builtin_template,
)
from backend.services import document_renderer, preparation_renderer


@pytest.fixture
def preparation_data():
    input_data = {
        "subject": "Python",
        "grade": "Year 1",
        "topic": "List operations",
        "duration": "2 hours",
        "location": "Lab 301",
        "textbook_name": "Python Basics",
    }
    content = {
        "teaching_goals": {
            "knowledge": ["Understand list creation"],
            "ability": ["Use common list operations"],
            "quality": ["Develop careful coding habits"],
        },
        "key_points": "Create and access a list",
        "difficult_points": "Index and slice operations",
        "teaching_steps": [
            {
                "stage": "Introduction",
                "duration": "10 min",
                "teacher_activity": "Present a practical task",
                "student_activity": "Observe and answer",
                "design_intent": "Activate prior knowledge",
            },
            {
                "stage": "Practice",
                "duration": "60 min",
                "teacher_activity": "Guide implementation",
                "student_activity": "Complete the coding task",
                "design_intent": "Learn by doing",
            },
        ],
        "homework": {
            "required": "Complete list exercises",
            "optional": "Build a contact list",
        },
        "blackboard_design": "Create - Access - Update",
    }
    return input_data, content


def test_builtin_template_is_fixed_and_valid():
    path = get_builtin_template_path()
    report = validate_builtin_template()

    assert BUILTIN_TEMPLATE_ID == "yunlin-standard"
    assert path.name == "yunlin_lesson_plan.docx"
    assert "backend/resources/templates" in path.as_posix()
    assert report["is_valid"] is True
    assert report["errors"] == []
    assert REQUIRED_TEMPLATE_FIELDS.issubset(set(report["variables"]))
    assert len(report["sha256"]) == 64


def test_other_template_id_is_rejected():
    with pytest.raises(ValueError, match="云林"):
        require_valid_builtin_template("uploaded-template")


def test_template_upload_and_office_routes_are_removed():
    paths = set(app.openapi()["paths"])

    assert "/api/templates/upload" not in paths
    assert not any("onlyoffice" in path.lower() for path in paths)
    assert "/api/templates/validation" in paths
    assert "/api/templates/validation/all" in paths
    assert "/api/preparation" in paths


def test_handout_and_presentation_render(tmp_path: pathlib.Path, preparation_data):
    input_data, content = preparation_data
    renderer = preparation_renderer.PreparationRenderer()
    renderer.output_dir = tmp_path

    handout_path = pathlib.Path(renderer.render_handout("test-id", input_data, content))
    presentation_path = pathlib.Path(
        renderer.render_presentation("test-id", input_data, content)
    )

    assert handout_path.is_file()
    assert presentation_path.is_file()
    assert "List operations" in "\n".join(
        paragraph.text for paragraph in Document(handout_path).paragraphs
    )
    handout_paragraphs = Document(handout_path).paragraphs
    assert not any(
        paragraph.text.startswith("\n") or paragraph.text.endswith("\n")
        for paragraph in handout_paragraphs
    )
    assert any(
        paragraph.text == "____________________________________________________________"
        and paragraph.paragraph_format.keep_together
        for paragraph in handout_paragraphs
    )
    presentation = Presentation(presentation_path)
    assert len(presentation.slides) >= 7
    assert any(
        "学习目标" in shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_presentation_dense_slide_uses_text_autofit():
    renderer = preparation_renderer.PreparationRenderer()
    presentation = Presentation()
    bullets = [
        "教师引导：" + "讲解配置、验证和故障排查步骤。" * 20 + "教师内容结束",
        "学习活动：" + "完成配置、记录证据并相互检查。" * 20 + "学生活动结束",
        "学习提示：" + "对照验收标准复核学习产出。" * 20 + "提示内容结束",
    ]

    renderer._add_bullet_slide(presentation, "课堂实践", bullets, eyebrow="30分钟")

    assert len(presentation.slides) >= 3
    rendered_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "教师内容结束" in rendered_text
    assert "学生活动结束" in rendered_text
    assert "提示内容结束" in rendered_text
    for slide in presentation.slides:
        body = next(
            shape.text_frame
            for shape in slide.shapes
            if hasattr(shape, "text_frame")
            and any(label in shape.text for label in ("教师引导：", "学习活动：", "学习提示："))
        )
        assert body.auto_size == pptx_text_enum.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        assert sum(len(paragraph.text) for paragraph in body.paragraphs) <= 520


def test_yunlin_lesson_plan_renders(tmp_path: pathlib.Path, preparation_data):
    input_data, content = preparation_data
    renderer = document_renderer.DocumentRenderer()
    renderer.output_dir = tmp_path
    render_data = {
        **input_data,
        "teaching_topic": input_data["topic"],
        "class_name": "Class A",
        "week_number": "1",
        "lesson_number": "1",
        "references": "Python Basics",
        "ideological_political": "Professional responsibility",
        **content,
    }

    output_path = pathlib.Path(
        renderer.render_lesson_plan(str(get_builtin_template_path()), render_data)
    )

    assert output_path.is_file()
    document = Document(output_path)
    all_text = "\n".join(
        [paragraph.text for paragraph in document.paragraphs]
        + [cell.text for table in document.tables for row in table.rows for cell in row.cells]
    )
    assert "List operations" in all_text
    assert "{{" not in all_text
    assert "授课学时：  2 学时" in all_text
    objective_cell = document.tables[0].rows[7].cells[0]
    assert len(objective_cell.paragraphs) == 1
    assert "\n\n" not in objective_cell.text
    for row_index in range(7, 12):
        assert document.tables[0].rows[row_index]._tr.trPr.find(qn("w:cantSplit")) is None
    assert document.tables[1].rows[4]._tr.trPr.find(qn("w:cantSplit")) is not None
    assert document.tables[1].rows[5]._tr.trPr.find(qn("w:trHeight")) is None
    assert document.tables[1].rows[6]._tr.trPr.find(qn("w:trHeight")) is None


def test_batch_lesson_plan_compacts_each_homepage(tmp_path: pathlib.Path, preparation_data):
    input_data, content = preparation_data
    renderer = document_renderer.DocumentRenderer()
    renderer.output_dir = tmp_path
    render_data = {
        **input_data,
        **content,
        "class_name": "Class 1",
        "week_number": "1",
        "lesson_number": "1",
        "references": "Python Basics",
        "ideological_political": "Work carefully and verify results",
    }

    output_path = pathlib.Path(
        renderer.render_lesson_plans_document(
            str(get_builtin_template_path()),
            [dict(render_data), {**render_data, "lesson_number": "2"}],
            course_name="Python",
            document_number=1,
            week_number=1,
        )
    )

    document = Document(output_path)
    assert len(document.tables) == 4
    for table_index in (0, 2):
        objective_cell = document.tables[table_index].rows[7].cells[0]
        assert len(objective_cell.paragraphs) == 1
        assert "\n\n" not in objective_cell.text

    body = document.element.body
    assert body[-2].tag == qn("w:tbl")
    empty_body_paragraphs = [
        element
        for element in body
        if element.tag == qn("w:p")
        and not any(
            (node.tag == qn("w:t") and node.text and node.text.strip())
            or node.tag == qn("w:br")
            for node in element.iter()
        )
    ]
    assert len(empty_body_paragraphs) == 2
