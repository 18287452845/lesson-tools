import json
from pathlib import Path

import pytest
from docx import Document
from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Inches

from backend.api import analytics, batch, course_archives, resources
from backend.models.schemas import (
    CourseArchiveCreate,
    CourseArchiveUpdate,
    TeachingResourceCreate,
    TeachingResourceUpdate,
)
from backend.services import (
    ai_metrics,
    batch_execution,
    course_archive_service,
    document_preview,
    teaching_resource_service,
)


@pytest.mark.asyncio
async def test_resource_library_crud_filters_and_prompt_context(test_db, monkeypatch):
    monkeypatch.setattr(teaching_resource_service, "db", test_db)
    first = await teaching_resource_service.create_resource({
        "title": "森林防火案例", "resource_type": "case", "subject": "信息安全",
        "grade": "2024级", "content": "分析监控系统告警并制定处置流程。",
        "source_url": "https://example.test/case", "tags": ["安全", "项目"],
    })
    second = await teaching_resource_service.create_resource({
        "title": "课堂互评量规", "resource_type": "rubric", "subject": "信息安全",
        "content": "按证据完整度与操作规范评分。", "tags": ["评价"],
    })
    assert first["tags"] == ["安全", "项目"]
    assert (await teaching_resource_service.get_resource(first["id"]))["title"] == "森林防火案例"

    items, total = await teaching_resource_service.list_resources(
        search="告警", resource_type="case", subject="信息安全", grade="2024级"
    )
    assert total == 1 and items[0]["id"] == first["id"]
    items, total = await teaching_resource_service.list_resources(page=1, limit=1)
    assert total == 2 and len(items) == 1

    updated = await teaching_resource_service.update_resource(
        first["id"], {"title": "森林防火应急案例", "tags": ["安全"], "ignored": 1}
    )
    assert updated["title"].endswith("案例") and updated["tags"] == ["安全"]
    assert await teaching_resource_service.update_resource("missing", {"title": "x"}) is None

    context = await teaching_resource_service.get_resource_context(
        [first["id"], second["id"], first["id"]], increment_use=True
    )
    assert "森林防火" in context and "课堂互评" in context
    assert (await teaching_resource_service.get_resource(first["id"]))["use_count"] == 1
    assert await teaching_resource_service.get_resource_context([]) == ""
    assert await teaching_resource_service.delete_resource(first["id"])
    assert not await teaching_resource_service.delete_resource("missing")
    archived, count = await teaching_resource_service.list_resources(status="archived")
    assert count == 1 and archived[0]["status"] == "archived"


@pytest.mark.asyncio
async def test_resource_api_success_and_not_found(test_db, monkeypatch):
    monkeypatch.setattr(teaching_resource_service, "db", test_db)
    created = await resources.create_resource(TeachingResourceCreate(
        title="实验任务", resource_type="experiment", content="完成部署与验收。", tags=["实训"]
    ))
    listing = await resources.list_resources(
        search=None, resource_type=None, subject=None, grade=None,
        status="active", page=1, limit=50,
    )
    assert listing.total == 1
    assert (await resources.get_resource(created["id"]))["id"] == created["id"]
    changed = await resources.update_resource(
        created["id"], TeachingResourceUpdate(content="完成部署、验证与复盘。")
    )
    assert "复盘" in changed["content"]
    await resources.delete_resource(created["id"])
    with pytest.raises(HTTPException) as exc:
        await resources.get_resource("missing")
    assert exc.value.status_code == 404


@pytest.mark.asyncio
async def test_course_archive_create_update_list_clone_and_archive(test_db, monkeypatch):
    monkeypatch.setattr(teaching_resource_service, "db", test_db)
    monkeypatch.setattr(course_archive_service, "db", test_db)
    resource = await teaching_resource_service.create_resource({
        "title": "项目任务", "resource_type": "assignment", "content": "提交项目报告。",
        "tags": [],
    })
    archive = await course_archive_service.create_archive({
        "course_name": "服务器安全", "subject": "信息安全", "grade": "2024级",
        "academic_year": "2026-2027", "semester": 1, "teacher_name": "李老师",
        "total_hours": 64, "hours_per_lesson": 2, "start_week": 2,
        "class_ids": ["class-1"], "resource_ids": [resource["id"], "missing"],
        "location": "实训楼", "notes": "秋季档案",
    })
    assert archive["resource_ids"] == [resource["id"]]
    assert archive["batch_task_count"] == archive["lesson_plan_count"] == 0
    items, total = await course_archive_service.list_archives(
        search="服务器", academic_year="2026-2027", semester=1
    )
    assert total == 1 and items[0]["id"] == archive["id"]

    changed = await course_archive_service.update_archive(
        archive["id"], {"teacher_name": "王老师", "class_ids": ["class-2"], "resource_ids": []}
    )
    assert changed["teacher_name"] == "王老师"
    assert changed["class_ids"] == ["class-2"] and changed["resource_ids"] == []
    assert await course_archive_service.update_archive("missing", {}) is None

    clone = await course_archive_service.clone_archive(archive["id"], "2027-2028", 2)
    assert clone["id"] != archive["id"] and clone["semester"] == 2
    assert await course_archive_service.clone_archive("missing", "2027-2028", 2) is None
    assert await course_archive_service.archive_course(archive["id"])
    assert not await course_archive_service.archive_course("missing")
    active, total = await course_archive_service.list_archives(limit=10)
    assert total == 1 and active[0]["id"] == clone["id"]


@pytest.mark.asyncio
async def test_course_archive_api_paths(test_db, monkeypatch):
    monkeypatch.setattr(course_archive_service, "db", test_db)
    request = CourseArchiveCreate(
        course_name="Python", subject="软件技术", grade="2025级",
        academic_year="2026-2027", semester=2, total_hours=32,
    )
    created = await course_archives.create_archive(request)
    listing = await course_archives.list_archives(
        search=None, academic_year="2026-2027", semester=None,
        status="active", page=1, limit=50,
    )
    assert listing.total == 1
    fetched = await course_archives.get_archive(created["id"])
    assert fetched["course_name"] == "Python"
    updated = await course_archives.update_archive(
        created["id"], CourseArchiveUpdate(location="机房 301")
    )
    assert updated["location"] == "机房 301"
    await course_archives.delete_archive(created["id"])
    with pytest.raises(HTTPException):
        await course_archives.update_archive("missing", CourseArchiveUpdate(notes="x"))


def test_docx_and_pptx_html_preview(tmp_path: Path):
    docx_path = tmp_path / "preview.docx"
    document = Document()
    document.add_heading("课程教案", level=0)
    document.add_heading("教学目标", level=1)
    paragraph = document.add_paragraph()
    paragraph.add_run("掌握基础").bold = True
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "课程"
    table.cell(0, 1).text = "Python"
    document.save(docx_path)
    html = document_preview.render_document_preview(docx_path)
    assert "课程教案" in html and "<table>" in html and "<strong>" in html

    pptx_path = tmp_path / "preview.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text_frame.paragraphs[0].text = "课堂演示"
    presentation.save(pptx_path)
    html = document_preview.render_document_preview(pptx_path)
    assert "课堂演示" in html and "aspect-ratio:16/9" in html
    with pytest.raises(ValueError):
        document_preview.render_document_preview(tmp_path / "other.pdf")


@pytest.mark.asyncio
async def test_ai_cost_quality_recording_and_analytics(test_db, monkeypatch):
    monkeypatch.setattr(ai_metrics, "db", test_db)
    monkeypatch.setattr(analytics, "db", test_db)
    monkeypatch.setattr(ai_metrics.os, "getenv", lambda _name: None)
    assert ai_metrics.estimate_cost("unknown", 1000, 1000) == 0
    assert ai_metrics.estimate_cost("deepseek-v4-flash", 1_000_000, 1_000_000) == 0.42
    await ai_metrics.record_ai_usage(
        provider="deepseek", model="deepseek-v4-flash", status="success",
        prompt_tokens=1000, cached_input_tokens=200, completion_tokens=500, latency_ms=800,
    )
    await ai_metrics.record_ai_usage(
        provider="deepseek", model="deepseek-v4-flash", status="failed",
        prompt_tokens=100, completion_tokens=0, latency_ms=100, error_message="timeout",
    )
    content = {
        "teaching_goals": {"knowledge": ["目标1", "目标2"], "ability": ["目标3"], "quality": ["目标4", "目标5"]},
        "key_points": "重点", "difficult_points": "难点", "homework": {"required": "任务"},
        "blackboard_design": "板书", "teaching_steps": [
            {"duration": "10分钟", "teacher_activity": "讲解" * 80, "student_activity": "实践" * 80}
            for _ in range(5)
        ],
    }
    score, dimensions = ai_metrics.evaluate_lesson_quality(content)
    assert score >= 90 and dimensions["completeness"] == 30
    empty_score, _ = ai_metrics.evaluate_lesson_quality({})
    assert empty_score == 0
    recorded = await ai_metrics.record_quality("preparation", "plan-1", content)
    assert recorded == score

    usage = await analytics.ai_summary(days=30)
    assert usage["summary"]["calls"] == 2 and usage["summary"]["success_rate"] == 50
    assert usage["by_model"][0]["model"] == "deepseek-v4-flash"
    quality = await analytics.quality_summary(days=30)
    assert quality["count"] == 1 and quality["average_score"] == score


@pytest.mark.asyncio
async def test_batch_launch_recovery_and_restart_checkpoint(test_db, monkeypatch):
    launched = []
    class Processor:
        def __init__(self, **kwargs): self.kwargs = kwargs
        async def process_batch_task(self, task_id, is_draft_mode=False): return None
    def background(coro, name):
        launched.append(name); coro.close()
    monkeypatch.setattr(batch_execution, "BatchTaskProcessor", Processor)
    monkeypatch.setattr(batch_execution, "run_in_background", background)
    batch_execution.launch_batch_task("normal-1", 2)
    batch_execution.launch_batch_task("draft-1", 2, "draft")
    assert launched == ["batch-task-normal-1", "draft-task-draft-1"]

    monkeypatch.setattr(batch_execution, "db", test_db)
    await test_db.execute(
        "INSERT INTO templates (id, name, file_path) VALUES ('t', 'T', 't.docx')", commit=True
    )
    await test_db.execute(
        """
        INSERT INTO batch_tasks
          (id, course_name, subject, grade, template_id, total_hours, hours_per_lesson,
           chapters, status, total_count, task_type)
        VALUES ('recover', 'C', 'S', 'G', 't', 2, 2, '[]', 'processing', 1, 'normal')
        """, commit=True,
    )
    recovered = []
    monkeypatch.setattr(batch_execution, "launch_batch_task", lambda *args: recovered.append(args))
    assert await batch_execution.recover_interrupted_batch_tasks() == 1
    assert recovered[0][0] == "recover"

    monkeypatch.setattr(batch, "get_db", lambda: _async_value(test_db))
    monkeypatch.setattr(batch, "launch_batch_task", lambda *args: recovered.append(args))
    await test_db.execute(
        "UPDATE batch_tasks SET status='failed', failed_count=1 WHERE id='recover'", commit=True
    )
    response = await batch._restart_checkpointed_task("recover", failed_only=True)
    assert response.status == "pending"
    with pytest.raises(HTTPException) as exc:
        await batch._restart_checkpointed_task("recover", failed_only=False)
    assert exc.value.status_code == 409
    with pytest.raises(HTTPException):
        await batch._restart_checkpointed_task("missing", failed_only=False)


async def _async_value(value):
    return value
