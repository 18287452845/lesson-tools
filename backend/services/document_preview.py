"""Safe, dependency-light HTML previews for generated DOCX and PPTX files."""
from __future__ import annotations

from html import escape
from pathlib import Path
from typing import Iterator, Union

from docx import Document
from docx.document import Document as DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from pptx import Presentation


def _page(title: str, body: str, extra_css: str = "") -> str:
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width">
<title>{escape(title)}</title><style>
*{{box-sizing:border-box}} body{{margin:0;background:#eef2f0;color:#20352f;font-family:'Microsoft YaHei','PingFang SC',sans-serif}}
.toolbar{{position:sticky;top:0;z-index:20;padding:10px 18px;background:#174a3d;color:white;box-shadow:0 2px 8px #0002}}
{extra_css}</style></head><body><div class="toolbar">只读预览 · {escape(title)}</div>{body}</body></html>"""


def _iter_docx_blocks(document: DocxDocument) -> Iterator[Union[Paragraph, Table]]:
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _paragraph_html(paragraph: Paragraph) -> str:
    pieces: list[str] = []
    for run in paragraph.runs:
        text = escape(run.text).replace("\n", "<br>")
        if not text:
            continue
        if run.bold:
            text = f"<strong>{text}</strong>"
        if run.italic:
            text = f"<em>{text}</em>"
        pieces.append(text)
    text = "".join(pieces) or escape(paragraph.text)
    if not text:
        return '<div class="spacer"></div>'
    style = (paragraph.style.name if paragraph.style else "").lower()
    if "title" in style:
        return f"<h1>{text}</h1>"
    if "heading 1" in style:
        return f"<h2>{text}</h2>"
    if "heading" in style:
        return f"<h3>{text}</h3>"
    list_class = ' class="list"' if "list" in style else ""
    return f"<p{list_class}>{text}</p>"


def docx_to_html(path: Path) -> str:
    document = Document(path)
    blocks: list[str] = []
    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            blocks.append(_paragraph_html(block))
            continue
        rows = []
        for row in block.rows:
            cells = "".join(
                f"<td>{escape(cell.text).replace(chr(10), '<br>')}</td>"
                for cell in row.cells
            )
            rows.append(f"<tr>{cells}</tr>")
        blocks.append(f"<table>{''.join(rows)}</table>")
    css = """
.paper{width:min(210mm,calc(100% - 28px));min-height:297mm;margin:22px auto;padding:20mm 18mm;background:white;box-shadow:0 5px 24px #183c3022}
h1{text-align:center;font-size:26px;margin:0 0 22px} h2{font-size:19px;color:#176b52;border-bottom:2px solid #d8e7e1;padding-bottom:6px} h3{font-size:16px;color:#285c4d}
p{font-family:'SimSun','Songti SC',serif;line-height:1.75;margin:8px 0;white-space:pre-wrap}.list:before{content:'• ';color:#176b52}.spacer{height:10px}
table{width:100%;border-collapse:collapse;margin:12px 0;font-size:14px}td{border:1px solid #718b82;padding:7px;line-height:1.5;vertical-align:top}
@media(max-width:700px){.paper{padding:20px 16px;margin:10px auto;min-height:auto}}
"""
    return _page(path.name, f'<main class="paper">{"".join(blocks)}</main>', css)


def _rgb(value, fallback: str) -> str:
    try:
        rgb = value.rgb
        return f"#{rgb}" if rgb is not None else fallback
    except (AttributeError, TypeError, ValueError):
        return fallback


def _contrast_text_color(background: str) -> str:
    """Choose a legible fallback for theme/inherited text colors."""
    try:
        red, green, blue = (
            int(background[index:index + 2], 16) / 255
            for index in (1, 3, 5)
        )
    except (TypeError, ValueError):
        return "#20352f"
    luminance = 0.2126 * red + 0.7152 * green + 0.0722 * blue
    return "#ffffff" if luminance < 0.45 else "#20352f"


def pptx_to_html(path: Path) -> str:
    presentation = Presentation(path)
    slide_width = float(presentation.slide_width)
    slide_height = float(presentation.slide_height)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        try:
            background = _rgb(slide.background.fill.fore_color, "#ffffff")
        except (AttributeError, TypeError, ValueError):
            background = "#ffffff"
        fallback_text_color = _contrast_text_color(background)
        shapes: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            left = float(shape.left) / slide_width * 100
            top = float(shape.top) / slide_height * 100
            width = float(shape.width) / slide_width * 100
            height = float(shape.height) / slide_height * 100
            paragraphs = []
            for paragraph in shape.text_frame.paragraphs:
                runs = []
                for run in paragraph.runs:
                    size = run.font.size.pt if run.font.size else 20
                    color = _rgb(run.font.color, fallback_text_color)
                    weight = "700" if run.font.bold else "400"
                    runs.append(
                        f'<span style="font-size:{max(10, min(size, 54))}px;color:{color};font-weight:{weight}">{escape(run.text)}</span>'
                    )
                paragraphs.append(f"<div>{''.join(runs) or escape(paragraph.text)}</div>")
            shapes.append(
                f'<div class="shape" style="left:{left:.2f}%;top:{top:.2f}%;width:{width:.2f}%;height:{height:.2f}%">{"".join(paragraphs)}</div>'
            )
        slides.append(
            f'<section><div class="slide" style="background:{background}">{"".join(shapes)}<span class="number">{index}</span></div></section>'
        )
    css = """
.deck{padding:24px}.deck section{max-width:1120px;margin:0 auto 28px}.slide{position:relative;width:100%;aspect-ratio:16/9;overflow:hidden;box-shadow:0 6px 26px #173f332b;border-radius:5px}
.shape{position:absolute;overflow:hidden;line-height:1.28;white-space:pre-wrap}.shape>div{margin:0 0 7px}.number{position:absolute;right:14px;bottom:9px;color:#61756e;font-size:12px}
@media(max-width:700px){.deck{padding:10px}.shape span{font-size:clamp(8px,2.5vw,18px)!important}}
"""
    return _page(path.name, f'<main class="deck">{"".join(slides)}</main>', css)


def render_document_preview(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return docx_to_html(path)
    if suffix == ".pptx":
        return pptx_to_html(path)
    raise ValueError("仅支持预览 DOCX 和 PPTX 文件")
