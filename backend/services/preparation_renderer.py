"""Render student handouts and classroom presentations from prepared content."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Iterable

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptPt

from ..config import settings


YUNLIN_GREEN = "176B52"
YUNLIN_DARK = PptColor(24, 74, 61)
YUNLIN_ACCENT = PptColor(36, 142, 108)


def _safe_filename(value: str) -> str:
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", value).strip(" .")
    return value[:60] or "备课资料"


def _as_lines(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [line.strip() for line in value.splitlines() if line.strip()]
    if isinstance(value, dict):
        result: list[str] = []
        for item in value.values():
            result.extend(_as_lines(item))
        return result
    if isinstance(value, Iterable):
        result = []
        for item in value:
            result.extend(_as_lines(item))
        return result
    return [str(value)]


class PreparationRenderer:
    """Create the non-template artifacts in the Yunlin visual style."""

    def __init__(self) -> None:
        self.output_dir = settings.output_dir

    def render_handout(
        self,
        preparation_id: str,
        input_data: dict[str, Any],
        content: dict[str, Any],
    ) -> str:
        document = Document()
        section = document.sections[0]
        section.top_margin = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin = Cm(2.0)
        section.right_margin = Cm(2.0)

        normal = document.styles["Normal"]
        normal.font.name = "Microsoft YaHei"
        normal._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        normal.font.size = Pt(10.5)

        title = document.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title.add_run(f"{input_data['topic']} · 学生讲义")
        title_run.bold = True
        title_run.font.size = Pt(22)
        title_run.font.color.rgb = RGBColor.from_string(YUNLIN_GREEN)
        title_run.font.name = "Microsoft YaHei"
        title_run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

        meta = document.add_table(rows=2, cols=3)
        meta.style = "Table Grid"
        meta_values = [
            ("课程", input_data.get("subject", "")),
            ("年级", input_data.get("grade", "")),
            ("课时", input_data.get("duration", "")),
            ("地点", input_data.get("location", "") or "—"),
            ("教材", input_data.get("textbook_name", "") or "—"),
            ("姓名", "________________"),
        ]
        for cell, (label, value) in zip(
            (cell for row in meta.rows for cell in row.cells), meta_values
        ):
            cell.text = f"{label}：{value}"

        self._add_docx_section(
            document,
            "一、学习目标",
            _as_lines(content.get("teaching_goals")),
        )
        self._add_docx_section(
            document,
            "二、学习重点与难点",
            [
                f"重点：{content.get('key_points', '')}",
                f"难点：{content.get('difficult_points', '')}",
            ],
        )

        heading = document.add_heading("三、课堂学习单", level=1)
        heading.runs[0].font.color.rgb = RGBColor.from_string(YUNLIN_GREEN)
        steps = content.get("teaching_steps") or []
        for index, step in enumerate(steps, start=1):
            if not isinstance(step, dict):
                continue
            stage = step.get("stage") or f"学习环节 {index}"
            duration = step.get("duration") or ""
            step_heading = document.add_heading(f"{index}. {stage} {duration}", level=2)
            step_heading.paragraph_format.keep_with_next = True
            activity = step.get("student_activity") or step.get("teacher_activity") or ""
            if activity:
                activity_paragraph = document.add_paragraph(str(activity))
                activity_paragraph.paragraph_format.keep_with_next = True
            note = document.add_paragraph("我的记录：")
            note.runs[0].bold = True
            note.paragraph_format.keep_with_next = True
            record_line = document.add_paragraph(
                "____________________________________________________________"
            )
            record_line.paragraph_format.keep_together = True

        homework = content.get("homework") or {}
        tasks = []
        if isinstance(homework, dict):
            if homework.get("required"):
                tasks.append(f"基础任务：{homework['required']}")
            if homework.get("optional"):
                tasks.append(f"拓展任务：{homework['optional']}")
        else:
            tasks = _as_lines(homework)
        self._add_docx_section(document, "四、课后巩固", tasks or ["根据课堂内容完成复习与实践。"])

        reflection = document.add_heading("五、学习反思", level=1)
        reflection.runs[0].font.color.rgb = RGBColor.from_string(YUNLIN_GREEN)
        document.add_paragraph(
            "本节课我掌握了：____________________________________________\n\n"
            "我仍需解决的问题：__________________________________________"
        )

        filename = f"{_safe_filename(input_data['topic'])}_学生讲义_{preparation_id[:8]}.docx"
        output_path = self.output_dir / filename
        document.save(output_path)
        return str(output_path)

    def _add_docx_section(
        self, document: Document, title: str, items: list[str]
    ) -> None:
        heading = document.add_heading(title, level=1)
        heading.runs[0].font.color.rgb = RGBColor.from_string(YUNLIN_GREEN)
        for item in items:
            if item:
                document.add_paragraph(item, style="List Bullet")

    def render_presentation(
        self,
        preparation_id: str,
        input_data: dict[str, Any],
        content: dict[str, Any],
    ) -> str:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        self._add_title_slide(
            presentation,
            input_data["topic"],
            f"{input_data.get('subject', '')} · {input_data.get('grade', '')} · {input_data.get('duration', '')}",
        )
        self._add_bullet_slide(
            presentation,
            "学习目标",
            _as_lines(content.get("teaching_goals")) or ["明确本课学习任务"],
        )
        self._add_bullet_slide(
            presentation,
            "重点与难点",
            [
                f"重点：{content.get('key_points', '')}",
                f"难点：{content.get('difficult_points', '')}",
            ],
        )

        steps = content.get("teaching_steps") or []
        for index, step in enumerate(steps[:8], start=1):
            if not isinstance(step, dict):
                continue
            stage = step.get("stage") or f"课堂环节 {index}"
            duration = step.get("duration") or ""
            bullets = []
            if step.get("teacher_activity"):
                bullets.append(f"教师引导：{step['teacher_activity']}")
            if step.get("student_activity"):
                bullets.append(f"学习活动：{step['student_activity']}")
            if step.get("design_intent"):
                bullets.append(f"学习提示：{step['design_intent']}")
            self._add_bullet_slide(
                presentation,
                f"{index:02d}  {stage}",
                bullets or ["围绕本环节任务开展学习活动"],
                eyebrow=duration,
            )

        summary_items = _as_lines(content.get("blackboard_design"))
        if not summary_items:
            summary_items = [
                str(content.get("key_points") or "回顾本课核心知识"),
                "完成自我检查，提出仍需解决的问题",
            ]
        self._add_bullet_slide(presentation, "课堂小结", summary_items)

        homework = content.get("homework") or {}
        homework_items = _as_lines(homework)
        self._add_bullet_slide(
            presentation,
            "课后任务",
            homework_items or ["完成课堂巩固任务并做好复习"],
        )

        filename = f"{_safe_filename(input_data['topic'])}_课堂PPT_{preparation_id[:8]}.pptx"
        output_path = self.output_dir / filename
        presentation.save(output_path)
        return str(output_path)

    def _add_title_slide(self, presentation: Presentation, title: str, subtitle: str) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = YUNLIN_DARK

        accent = slide.shapes.add_shape(1, Inches(0.75), Inches(0.8), Inches(0.16), Inches(4.8))
        accent.fill.solid()
        accent.fill.fore_color.rgb = YUNLIN_ACCENT
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1.25), Inches(1.4), Inches(10.8), Inches(2.2))
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = PptPt(50)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptColor(255, 255, 255)

        sub_box = slide.shapes.add_textbox(Inches(1.3), Inches(4.1), Inches(10), Inches(0.8))
        sub = sub_box.text_frame.paragraphs[0]
        sub.text = subtitle
        sub.font.name = "Microsoft YaHei"
        sub.font.size = PptPt(24)
        sub.font.color.rgb = PptColor(193, 224, 211)

        footer = slide.shapes.add_textbox(Inches(1.3), Inches(6.35), Inches(10), Inches(0.4))
        footer_p = footer.text_frame.paragraphs[0]
        footer_p.text = "云南林业职业技术学院 · 智能备课"
        footer_p.font.name = "Microsoft YaHei"
        footer_p.font.size = PptPt(10)
        footer_p.font.color.rgb = PptColor(139, 184, 165)

    def _add_bullet_slide(
        self,
        presentation: Presentation,
        title: str,
        bullets: list[str],
        eyebrow: str = "",
    ) -> None:
        rendered_bullets = [
            chunk
            for bullet in bullets[:7]
            for chunk in self._split_presentation_bullet(str(bullet))
        ]
        pages: list[list[str]] = []
        current_page: list[str] = []
        current_characters = 0
        for bullet in rendered_bullets:
            if current_page and current_characters + len(bullet) > 520:
                pages.append(current_page)
                current_page = []
                current_characters = 0
            current_page.append(bullet)
            current_characters += len(bullet)
        if current_page or not pages:
            pages.append(current_page or ["围绕本环节任务开展学习活动"])

        for page_index, page_bullets in enumerate(pages):
            page_title = title if page_index == 0 else f"{title}（续）"
            self._add_bullet_slide_page(
                presentation,
                page_title,
                page_bullets,
                eyebrow,
            )

    @staticmethod
    def _split_presentation_bullet(text: str, max_characters: int = 480) -> list[str]:
        if len(text) <= max_characters:
            return [text]

        label = next(
            (
                candidate
                for candidate in (
                    "教师引导：",
                    "学习活动：",
                    "学习提示：",
                    "重点：",
                    "难点：",
                )
                if text.startswith(candidate)
            ),
            "",
        )
        remaining = text[len(label) :]
        chunks: list[str] = []
        prefix = label
        continuation_prefix = f"{label[:-1]}（续）：" if label else "续："
        while remaining:
            available = max_characters - len(prefix)
            if len(remaining) <= available:
                split_at = len(remaining)
            else:
                candidate = remaining[:available]
                split_at = max(candidate.rfind(mark) + 1 for mark in "。；！？")
                if split_at < available // 2:
                    split_at = available
            chunk = remaining[:split_at].strip()
            if chunk:
                chunks.append(prefix + chunk)
            remaining = remaining[split_at:].lstrip()
            prefix = continuation_prefix
        return chunks

    def _add_bullet_slide_page(
        self,
        presentation: Presentation,
        title: str,
        bullets: list[str],
        eyebrow: str = "",
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptColor(248, 251, 249)

        top_bar = slide.shapes.add_shape(1, 0, 0, presentation.slide_width, Inches(0.16))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = YUNLIN_ACCENT
        top_bar.line.fill.background()

        if eyebrow:
            eye_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(2.0), Inches(0.35))
            eye = eye_box.text_frame.paragraphs[0]
            eye.text = eyebrow
            eye.font.name = "Microsoft YaHei"
            eye.font.size = PptPt(11)
            eye.font.bold = True
            eye.font.color.rgb = YUNLIN_ACCENT

        title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.92), Inches(11.6), Inches(0.75))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.name = "Microsoft YaHei"
        title_p.font.size = PptPt(35)
        title_p.font.bold = True
        title_p.font.color.rgb = YUNLIN_DARK

        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.1), Inches(4.7))
        frame = body_box.text_frame
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        character_count = sum(len(bullet) for bullet in bullets)
        if character_count > 760:
            font_size = 13
        elif character_count > 560:
            font_size = 15
        elif character_count > 420:
            font_size = 17
        else:
            font_size = 19 if len(bullets) <= 4 else 16
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = PptPt(font_size)
            paragraph.font.color.rgb = PptColor(44, 62, 55)
            paragraph.space_after = PptPt(13)
            paragraph.alignment = PP_ALIGN.LEFT

        page = slide.shapes.add_textbox(Inches(11.8), Inches(6.85), Inches(0.7), Inches(0.3))
        page_p = page.text_frame.paragraphs[0]
        page_p.text = f"{len(presentation.slides):02d}"
        page_p.font.size = PptPt(10)
        page_p.font.color.rgb = YUNLIN_ACCENT
        page_p.alignment = PP_ALIGN.RIGHT
